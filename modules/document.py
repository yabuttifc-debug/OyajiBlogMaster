from pathlib import Path


def extract_text(file_path: str) -> str:
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"ファイルが見つかりません: {file_path}")

    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _read_pdf(path)
    elif suffix == ".docx":
        return _read_docx(path)
    elif suffix in (".pptx", ".ppt"):
        return _read_pptx(path)
    elif suffix == ".txt":
        return path.read_text(encoding="utf-8")
    else:
        raise ValueError(f"対応していないファイル形式です: {suffix}\n対応形式: .pdf, .docx, .pptx, .txt")


def _read_pdf(path: Path) -> str:
    import pdfplumber
    texts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
    return "\n\n".join(texts)


def _read_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def _read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            texts.append(f"【スライド {i}】\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)
